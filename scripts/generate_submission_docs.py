from __future__ import annotations

import json
import re
from collections import Counter
from dataclasses import dataclass
from datetime import datetime, timezone
from pathlib import Path

import matplotlib.pyplot as plt
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches, Pt
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PInches
from pptx.util import Pt as PPt


ROOT = Path(__file__).resolve().parents[1]
OUTPUTS = ROOT / "outputs"
DOCS = ROOT / "docs"
SUBMISSIONS = ROOT / "submissions"
SUBMISSIONS.mkdir(exist_ok=True)
CHARTS_DIR = SUBMISSIONS / "assets"
CHARTS_DIR.mkdir(exist_ok=True)

REPORT_JSON = OUTPUTS / "report.json"
BASELINE_JSON = OUTPUTS / "baseline.semgrep.json"
COMPARISON_MD = OUTPUTS / "comparison.md"
PROMPT_LOG = ROOT / "prompt_log.md"

STUDENT_NAME = "Charbel Lebbous"
STUDENT_ID = "8324"
COURSE = "CBRS503"
TITLE = "Secure Code Inspector"


@dataclass
class VersionResult:
    version: str
    title: str
    changes: list[str]
    why: list[str]
    ai_verified: int
    semgrep_verified: int
    matched: int
    precision_pct: float
    recall_pct: float
    ai_only: int
    semgrep_only: int
    interpretation: list[str]


def _extract_list(text: str, section_name: str) -> list[str]:
    pattern = rf"- {re.escape(section_name)}:\n((?:\s+- .+\n)+)"
    match = re.search(pattern, text)
    if not match:
        return []
    block = match.group(1)
    return [line.strip()[2:].strip() for line in block.splitlines() if line.strip().startswith("- ")]


def parse_prompt_versions(prompt_text: str) -> list[VersionResult]:
    sections = re.split(r"\n## (v\d+) - ", prompt_text)[1:]
    versions: list[VersionResult] = []
    for i in range(0, len(sections), 2):
        version = sections[i].strip()
        body = sections[i + 1]
        title, _, content = body.partition("\n")
        changes = _extract_list(content, "What changed")
        if not changes:
            changes = _extract_list(content, "Prompt/Config profile")
        why = _extract_list(content, "Why this improved")
        if not why:
            why = _extract_list(content, "Why this should improve")
        interpretation = _extract_list(content, "Interpretation")

        def extract_int(label: str) -> int:
            m = re.search(rf"- {re.escape(label)}:\s*(\d+)", content)
            return int(m.group(1)) if m else 0

        def extract_pct(label: str) -> float:
            m = re.search(rf"- {re.escape(label)}:\s*\*\*(\d+(?:\.\d+)?)%\*\*", content)
            return float(m.group(1)) if m else 0.0

        versions.append(
            VersionResult(
                version=version,
                title=title.strip(),
                changes=changes,
                why=why,
                ai_verified=extract_int("AI verified"),
                semgrep_verified=extract_int("Semgrep verified"),
                matched=extract_int("Matched"),
                precision_pct=extract_pct("Precision"),
                recall_pct=extract_pct("Recall"),
                ai_only=extract_int("AI-only"),
                semgrep_only=extract_int("Semgrep-only"),
                interpretation=interpretation,
            )
        )
    return versions


def parse_comparison_metrics(comparison_md: str) -> dict:
    row_match = re.search(
        r"\|\s*Semgrep\s*\|\s*(\d+)\s*\|\s*(\d+)\s*\|\s*(\d+)\s*\|\s*(\d+)\s*\|\s*(\d+)\s*\|\s*([0-9.]+)\s*\|\s*([0-9.]+)\s*\|",
        comparison_md,
    )
    metrics = {
        "ai_verified": 0,
        "semgrep_verified": 0,
        "matched": 0,
        "ai_only": 0,
        "semgrep_only": 0,
        "precision": 0.0,
        "recall": 0.0,
    }
    if row_match:
        metrics = {
            "ai_verified": int(row_match.group(1)),
            "semgrep_verified": int(row_match.group(2)),
            "matched": int(row_match.group(3)),
            "ai_only": int(row_match.group(4)),
            "semgrep_only": int(row_match.group(5)),
            "precision": float(row_match.group(6)),
            "recall": float(row_match.group(7)),
        }

    def section_items(name: str) -> list[str]:
        pat = rf"## {re.escape(name)}\n\n((?:- .+\n)+)"
        m = re.search(pat, comparison_md)
        if not m:
            return []
        return [ln.strip()[2:].strip() for ln in m.group(1).splitlines() if ln.strip().startswith("- ")]

    metrics["analysis_summary"] = section_items("Analysis Summary")
    metrics["false_positives"] = section_items("False Positives (AI)")
    metrics["misses"] = section_items("Misses (AI)")
    metrics["why_fp"] = section_items("Why False Positives Happen")
    metrics["why_miss"] = section_items("Why Misses Happen")
    return metrics


def build_charts(versions: list[VersionResult], cmp_metrics: dict) -> dict[str, Path]:
    chart_paths: dict[str, Path] = {}

    ver_names = [v.version for v in versions]
    precisions = [v.precision_pct for v in versions]
    recalls = [v.recall_pct for v in versions]
    plt.figure(figsize=(8, 4.2))
    plt.plot(ver_names, precisions, marker="o", linewidth=2.2, label="Precision %")
    plt.plot(ver_names, recalls, marker="o", linewidth=2.2, label="Recall %")
    plt.title("Prompt Version Progress (v1 to v7)")
    plt.ylabel("Percentage")
    plt.ylim(0, max(70, max(precisions + recalls) + 5))
    plt.grid(alpha=0.3)
    plt.legend()
    progress_png = CHARTS_DIR / "version_progress.png"
    plt.tight_layout()
    plt.savefig(progress_png, dpi=180)
    plt.close()
    chart_paths["version_progress"] = progress_png

    labels = ["Matched", "AI-only", "Semgrep-only"]
    values = [cmp_metrics["matched"], cmp_metrics["ai_only"], cmp_metrics["semgrep_only"]]
    plt.figure(figsize=(6, 4.2))
    plt.pie(values, labels=labels, autopct="%1.0f%%", startangle=120)
    plt.title("Current Comparison Distribution")
    pie_png = CHARTS_DIR / "comparison_distribution.png"
    plt.tight_layout()
    plt.savefig(pie_png, dpi=180)
    plt.close()
    chart_paths["comparison_distribution"] = pie_png

    return chart_paths


def add_toc(doc: Document) -> None:
    p = doc.add_paragraph()
    run = p.add_run()
    fld_char_begin = OxmlElement("w:fldChar")
    fld_char_begin.set(qn("w:fldCharType"), "begin")
    instr_text = OxmlElement("w:instrText")
    instr_text.set(qn("xml:space"), "preserve")
    instr_text.text = r'TOC \o "1-3" \h \z \u'
    fld_char_separate = OxmlElement("w:fldChar")
    fld_char_separate.set(qn("w:fldCharType"), "separate")
    fld_char_end = OxmlElement("w:fldChar")
    fld_char_end.set(qn("w:fldCharType"), "end")
    run._r.append(fld_char_begin)
    run._r.append(instr_text)
    run._r.append(fld_char_separate)
    run._r.append(fld_char_end)


def add_bullets(doc: Document, items: list[str]) -> None:
    for item in items:
        p = doc.add_paragraph(style="List Bullet")
        p.add_run(item)


def build_report_docx(
    report: dict,
    baseline: dict,
    comparison_md: str,
    versions: list[VersionResult],
    cmp_metrics: dict,
    chart_paths: dict[str, Path],
) -> Path:
    now = datetime.now(timezone.utc).strftime("%Y-%m-%d %H:%M UTC")
    findings = report.get("findings", [])
    cats = Counter(f.get("owasp_category", "Unknown") for f in findings)
    agents = Counter(f.get("source_agent", "Unknown") for f in findings)
    run_meta = report.get("run_metadata", {})

    doc = Document()
    doc.styles["Normal"].font.name = "Calibri"
    doc.styles["Normal"].font.size = Pt(11)

    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    tr = title.add_run(f"{COURSE} Final Project Report\n{TITLE}")
    tr.bold = True
    tr.font.size = Pt(24)
    meta = doc.add_paragraph()
    meta.alignment = WD_ALIGN_PARAGRAPH.CENTER
    meta.add_run(f"Student: {STUDENT_NAME}\nStudent ID: {STUDENT_ID}\nGenerated: {now}")

    doc.add_page_break()
    doc.add_heading("Table of Contents", level=1)
    add_toc(doc)
    doc.add_paragraph("Note: In Word, right-click the table and select 'Update Field' to refresh page numbers.")
    doc.add_page_break()

    doc.add_heading("1. Executive Summary", level=1)
    doc.add_paragraph(
        "This project implements a web-first AI secure code inspector for OWASP-oriented vulnerability analysis. "
        "The tool scans a fixed repository scope, applies a configurable multi-agent prompt workflow, generates "
        "structured findings, and compares AI findings against a Semgrep baseline with precision/recall metrics and analysis."
    )

    doc.add_heading("2. PDF Requirements Compliance Matrix", level=1)
    tbl = doc.add_table(rows=1, cols=4)
    tbl.style = "Light List Accent 1"
    headers = ["Requirement from PDF", "Status", "Implemented How", "Evidence in Repo"]
    for i, h in enumerate(headers):
        tbl.rows[0].cells[i].text = h

    requirement_rows = [
        ("Tool format: CLI or lightweight web UI", "Done", "Streamlit web UI + CLI fallback", "web_app.py, src/secure_inspector/cli.py"),
        ("Input: repo path or selected files", "Done", "ZIP upload in UI, scoped file selection", "web_app.py, configs/scope.juiceshop.yaml"),
        ("Output: report.json + report.md", "Done", "Generated by Aggregator/reporting", "outputs/report.json, outputs/report.md"),
        ("Finding includes file + line range", "Done", "Finding schema enforces file and line spans", "src/secure_inspector/models.py"),
        ("OWASP category per finding", "Done", "Category mapping in prompts/data", "data/owasp_top10.yaml"),
        ("Risk summary (2-3 lines)", "Done", "Prompt constraints + report fields", "prompts/agents/*.md"),
        ("Fix recommendation (specific)", "Done", "Prompt constraints require actionable remediations", "report.json fix_recommendation"),
        ("Confidence score (0-1)", "Done", "Schema + verifier thresholding", "configs/profile.yaml, models.py"),
        ("Chunking + aggregation", "Done", "Chunker and Aggregator stage", "src/secure_inspector/chunker.py, agents/aggregator.py"),
        ("At least 3 prompt strategies + justification", "Done", "Role/constraints, few-shot, verification, OWASP injection, chunking", "prompts/, prompt_log.md"),
        ("Prompt log with versions and improvements", "Done", "v1 to v7 tracked with outcomes", "prompt_log.md"),
        ("Baseline comparison mandatory", "Done", "Semgrep baseline integrated", "src/secure_inspector/baseline/semgrep_runner.py"),
        ("Evaluation metrics + false positives + misses", "Done", "comparison.md analysis sections", "outputs/comparison.md"),
        ("GitHub reproducibility: code+prompts+outputs", "Done", "Repo structured with configs/prompts/outputs/docs", "README.md + folder structure"),
        ("One-command style run preferred", "Done", "python -m entrypoints documented", "README.md"),
        ("Fixed scope documented", "Done", "Juice Shop scope config", "configs/scope.juiceshop.yaml"),
        ("Demo link in README", "Pending final link", "Placeholder section kept for final add", "README.md"),
    ]
    for row in requirement_rows:
        cells = tbl.add_row().cells
        for i, val in enumerate(row):
            cells[i].text = val

    doc.add_heading("3. Final Architecture and Pipeline", level=1)
    doc.add_paragraph("4+1 orchestration implemented:")
    add_bullets(
        doc,
        [
            "Always-on agents: InjectionSpecialistAgent, AccessControlSpecialistAgent, VerifierAgent, AggregatorAgent.",
            "Conditional agent: ExtraCategorySpecialistAgent (enabled only when non-core categories exist in profile).",
            "Pipeline: ScopeLoader -> Chunker -> Specialist agents (parallel) -> Verifier -> Aggregator -> report artifacts.",
            "Baseline path: Semgrep on same scope -> normalized baseline artifact.",
            "Comparison path: AI-assisted semantic matching -> precision/recall + misses + false positives analysis.",
        ],
    )

    doc.add_heading("4. Prompt Engineering Design (What and Why)", level=1)
    ptable = doc.add_table(rows=1, cols=3)
    ptable.style = "Light Grid Accent 1"
    ptable.rows[0].cells[0].text = "Strategy"
    ptable.rows[0].cells[1].text = "Implementation"
    ptable.rows[0].cells[2].text = "Reason"
    prompt_rows = [
        ("Role + constraints", "Each agent prompt defines role scope, strict JSON output, abstain behavior.", "Reduce hallucinations and enforce deterministic structure."),
        ("Few-shot examples", "Positive and negative labeled examples in prompts/few_shot_examples.json.", "Improve response consistency and reduce false positives."),
        ("Verification pass", "VerifierAgent checks evidence, line validity, category consistency.", "Filter weak/speculative findings before final report."),
        ("OWASP definitions injection", "OWASP references provided to prompts through data and templates.", "Keep category mapping aligned with standard taxonomy."),
        ("Chunking strategy", "Code chunking with stable file/line mapping before LLM calls.", "Maintain traceability and avoid context overflow."),
        ("Final aggregation", "Aggregator merges/deduplicates/ranks findings and writes outputs.", "Produce clean final deliverables for grading."),
    ]
    for row in prompt_rows:
        c = ptable.add_row().cells
        c[0].text, c[1].text, c[2].text = row

    doc.add_heading("5. Prompt Version Evolution (v1 to v7)", level=1)
    summary = doc.add_table(rows=1, cols=8)
    summary.style = "Light List Accent 2"
    heads = ["Ver", "AI", "Semgrep", "Matched", "Precision", "Recall", "AI-only", "Semgrep-only"]
    for i, h in enumerate(heads):
        summary.rows[0].cells[i].text = h
    for v in versions:
        row = summary.add_row().cells
        row[0].text = v.version
        row[1].text = str(v.ai_verified)
        row[2].text = str(v.semgrep_verified)
        row[3].text = str(v.matched)
        row[4].text = f"{v.precision_pct:.2f}%"
        row[5].text = f"{v.recall_pct:.2f}%"
        row[6].text = str(v.ai_only)
        row[7].text = str(v.semgrep_only)

    if chart_paths.get("version_progress"):
        doc.add_paragraph("Trend chart:")
        doc.add_picture(str(chart_paths["version_progress"]), width=Inches(6.4))

    for v in versions:
        doc.add_heading(f"5.{v.version[1:]} {v.version} - {v.title}", level=2)
        doc.add_paragraph("What was changed:")
        add_bullets(doc, v.changes if v.changes else ["No explicit change list recorded."])
        if v.why:
            doc.add_paragraph("Why this change:")
            add_bullets(doc, v.why)
        doc.add_paragraph(
            f"Results: AI={v.ai_verified}, Baseline={v.semgrep_verified}, Matched={v.matched}, "
            f"Precision={v.precision_pct:.2f}%, Recall={v.recall_pct:.2f}%."
        )
        if v.interpretation:
            doc.add_paragraph("Analysis:")
            add_bullets(doc, v.interpretation)

    doc.add_heading("6. Current Run Results and Artifact Review", level=1)
    doc.add_paragraph(
        f"Run timestamp: {run_meta.get('timestamp', 'n/a')}\n"
        f"Model: {run_meta.get('model', 'n/a')}\n"
        f"Scope files: {len(run_meta.get('scope_files', []))}\n"
        f"Enabled OWASP categories: {', '.join(run_meta.get('enabled_categories', []))}"
    )
    ctable = doc.add_table(rows=1, cols=2)
    ctable.style = "Light List Accent 3"
    ctable.rows[0].cells[0].text = "OWASP Category"
    ctable.rows[0].cells[1].text = "AI Finding Count"
    for cat, count in sorted(cats.items()):
        row = ctable.add_row().cells
        row[0].text = cat
        row[1].text = str(count)

    atable = doc.add_table(rows=1, cols=2)
    atable.style = "Light List Accent 4"
    atable.rows[0].cells[0].text = "Source Agent"
    atable.rows[0].cells[1].text = "Findings"
    for ag, count in sorted(agents.items()):
        row = atable.add_row().cells
        row[0].text = ag
        row[1].text = str(count)

    doc.add_heading("7. Baseline Comparison and Analysis", level=1)
    metrics_tbl = doc.add_table(rows=1, cols=2)
    metrics_tbl.style = "Light Grid Accent 2"
    metrics_tbl.rows[0].cells[0].text = "Metric"
    metrics_tbl.rows[0].cells[1].text = "Value"
    rows = [
        ("AI verified findings", str(cmp_metrics["ai_verified"])),
        ("Semgrep verified findings", str(cmp_metrics["semgrep_verified"])),
        ("Matched findings", str(cmp_metrics["matched"])),
        ("AI-only findings", str(cmp_metrics["ai_only"])),
        ("Semgrep-only findings", str(cmp_metrics["semgrep_only"])),
        ("Precision", f"{cmp_metrics['precision'] * 100:.2f}%"),
        ("Recall", f"{cmp_metrics['recall'] * 100:.2f}%"),
    ]
    for key, val in rows:
        r = metrics_tbl.add_row().cells
        r[0].text = key
        r[1].text = val

    if chart_paths.get("comparison_distribution"):
        doc.add_paragraph("Distribution chart:")
        doc.add_picture(str(chart_paths["comparison_distribution"]), width=Inches(5.8))

    doc.add_paragraph("Comparison summary analysis:")
    add_bullets(doc, cmp_metrics.get("analysis_summary", []))

    doc.add_paragraph("False positives analysis (AI-only):")
    add_bullets(doc, cmp_metrics.get("false_positives", []))
    doc.add_paragraph("Misses analysis (Semgrep-only):")
    add_bullets(doc, cmp_metrics.get("misses", []))
    doc.add_paragraph("Root-cause explanation:")
    add_bullets(doc, cmp_metrics.get("why_fp", []))
    add_bullets(doc, cmp_metrics.get("why_miss", []))

    doc.add_heading("8. Reproducibility Checklist", level=1)
    add_bullets(
        doc,
        [
            "Source code, prompts, and configs are in repository structure.",
            "Fixed scope declared in configs/scope.juiceshop.yaml.",
            "One-command module entrypoints documented in README.",
            "Output artifacts generated in outputs/.",
            "Baseline output generated in outputs/baseline.semgrep.json.",
            "Comparison analysis available in outputs/comparison.md.",
            "Prompt log maintained in prompt_log.md.",
        ],
    )

    doc.add_heading("9. UI Evidence", level=1)
    images = [
        ("Run AI tab", "ui_run_ai_results.png"),
        ("Run Baseline tab", "ui_run_baseline_results.png"),
        ("Compare tab", "ui_compare_metrics.png"),
        ("Artifacts tab", "ui_artifacts_downloads.png"),
        ("Settings / Help tab", "ui_settings_help.png"),
    ]
    for caption, fname in images:
        path = DOCS / fname
        if path.exists():
            doc.add_paragraph(caption)
            doc.add_picture(str(path), width=Inches(6.5))

    doc.add_heading("10. Limitations and Improvement Plan", level=1)
    add_bullets(
        doc,
        [
            "Recall is still below the target threshold in some runs due to server.ts mismatch concentration.",
            "Category overlap between A01/A07/A02 can still produce semantic disagreement with baseline.",
            "Future pass: stronger server.ts-focused few-shots and verifier calibration per category.",
            "Optional: add second baseline (SonarQube/CodeQL) for triangulated evaluation.",
        ],
    )

    doc.add_heading("11. Conclusion", level=1)
    doc.add_paragraph(
        "The project satisfies the core CBRS503 requirements: configurable OWASP-focused AI analysis, prompt-engineering "
        "evidence with version log, mandatory baseline comparison, reproducibility, and structured artifacts ready for evaluation."
    )

    out = SUBMISSIONS / "CBRS503_Report_Charbel_Lebbous_8324.docx"
    doc.save(out)
    return out


def build_presentation(versions: list[VersionResult], cmp_metrics: dict, charts: dict[str, Path]) -> Path:
    prs = Presentation()
    prs.slide_width = PInches(13.333)
    prs.slide_height = PInches(7.5)

    def styled_title(slide, title_text: str, subtitle: str = ""):
        title_box = slide.shapes.add_textbox(PInches(0.6), PInches(0.4), PInches(12.2), PInches(1.0))
        tf = title_box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.bold = True
        p.font.size = PPt(34)
        p.font.color.rgb = RGBColor(25, 45, 75)
        if subtitle:
            sp = tf.add_paragraph()
            sp.text = subtitle
            sp.font.size = PPt(16)
            sp.font.color.rgb = RGBColor(70, 85, 100)

    # Slide 1
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    styled_title(s1, f"{COURSE} - {TITLE}", f"{STUDENT_NAME} | ID {STUDENT_ID}")
    b1 = s1.shapes.add_textbox(PInches(0.8), PInches(1.8), PInches(12.0), PInches(4.8)).text_frame
    bullets = [
        "Web-first AI secure code scanner",
        "OWASP-mapped findings + actionable fixes",
        "Semgrep baseline + AI-assisted comparison",
        "Reproducible pipeline and artifacts",
    ]
    for i, t in enumerate(bullets):
        p = b1.paragraphs[0] if i == 0 else b1.add_paragraph()
        p.text = t
        p.font.size = PPt(30 if i == 0 else 24)

    # Slide 2
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    styled_title(s2, "Scope and Deliverables")
    box = s2.shapes.add_textbox(PInches(0.8), PInches(1.5), PInches(12.0), PInches(5.2)).text_frame
    items = [
        "Input: Repository ZIP (fixed scope)",
        "Output: report.json, report.md",
        "Baseline: baseline.semgrep.json",
        "Evaluation: comparison.md (metrics + analysis)",
        "Config-driven categories: A03, A01, optional extras",
    ]
    for i, t in enumerate(items):
        p = box.paragraphs[0] if i == 0 else box.add_paragraph()
        p.text = t
        p.font.size = PPt(24)

    # Slide 3
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    styled_title(s3, "4+1 Agent Orchestration")
    box3 = s3.shapes.add_textbox(PInches(0.8), PInches(1.5), PInches(12.0), PInches(5.2)).text_frame
    for i, t in enumerate(
        [
            "ScopeLoader -> Chunker",
            "InjectionSpecialist + AccessControlSpecialist",
            "ExtraCategorySpecialist (conditional)",
            "VerifierAgent -> AggregatorAgent",
            "Outputs + baseline comparison",
        ]
    ):
        p = box3.paragraphs[0] if i == 0 else box3.add_paragraph()
        p.text = t
        p.font.size = PPt(25)

    # Slide 4
    s4 = prs.slides.add_slide(prs.slide_layouts[6])
    styled_title(s4, "Prompt Engineering (Implemented)")
    b4 = s4.shapes.add_textbox(PInches(0.8), PInches(1.5), PInches(12.0), PInches(5.2)).text_frame
    for i, t in enumerate(
        [
            "Role + constraints",
            "Few-shot examples",
            "Verification pass",
            "OWASP definitions injection",
            "Chunking + final aggregation",
            "Prompt Log: v1 to v7 with measured impact",
        ]
    ):
        p = b4.paragraphs[0] if i == 0 else b4.add_paragraph()
        p.text = t
        p.font.size = PPt(24)

    # Slide 5
    s5 = prs.slides.add_slide(prs.slide_layouts[6])
    styled_title(s5, "Results Snapshot")
    chart_data = CategoryChartData()
    chart_data.categories = ["Precision", "Recall"]
    chart_data.add_series("Current", [cmp_metrics["precision"] * 100, cmp_metrics["recall"] * 100])
    s5.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        PInches(0.8),
        PInches(1.5),
        PInches(6.0),
        PInches(4.3),
        chart_data,
    )
    t5 = s5.shapes.add_textbox(PInches(7.1), PInches(1.7), PInches(5.5), PInches(4.5)).text_frame
    lines = [
        f"AI verified: {cmp_metrics['ai_verified']}",
        f"Semgrep verified: {cmp_metrics['semgrep_verified']}",
        f"Matched: {cmp_metrics['matched']}",
        f"AI-only: {cmp_metrics['ai_only']}",
        f"Semgrep-only: {cmp_metrics['semgrep_only']}",
        f"Precision: {cmp_metrics['precision']*100:.0f}%",
        f"Recall: {cmp_metrics['recall']*100:.0f}%",
    ]
    for i, t in enumerate(lines):
        p = t5.paragraphs[0] if i == 0 else t5.add_paragraph()
        p.text = t
        p.font.size = PPt(21)

    # Slide 6
    s6 = prs.slides.add_slide(prs.slide_layouts[6])
    styled_title(s6, "Version Progress (v1 -> v7)")
    if charts.get("version_progress") and charts["version_progress"].exists():
        s6.shapes.add_picture(str(charts["version_progress"]), PInches(0.9), PInches(1.5), width=PInches(8.3))
    latest = versions[-1]
    info = s6.shapes.add_textbox(PInches(9.4), PInches(1.6), PInches(3.6), PInches(4.5)).text_frame
    points = [
        "v1: broad, noisy",
        "v4: higher coverage",
        "v5: precision guard",
        "v7: dedupe + semantic scoring",
        f"Final precision {latest.precision_pct:.0f}%",
        f"Final recall {latest.recall_pct:.0f}%",
    ]
    for i, t in enumerate(points):
        p = info.paragraphs[0] if i == 0 else info.add_paragraph()
        p.text = t
        p.font.size = PPt(18)

    # Slide 7
    s7 = prs.slides.add_slide(prs.slide_layouts[6])
    styled_title(s7, "Status and Next Technical Steps")
    b7 = s7.shapes.add_textbox(PInches(0.8), PInches(1.5), PInches(12.0), PInches(5.2)).text_frame
    for i, t in enumerate(
        [
            "Requirement coverage: complete except adding final demo link",
            "Strength: reproducible pipeline and clear artifacts",
            "Gap: recall still below target in server.ts-heavy cases",
            "Next: targeted server.ts few-shots + verifier calibration",
            "Deployment: public GitHub + Streamlit Cloud",
        ]
    ):
        p = b7.paragraphs[0] if i == 0 else b7.add_paragraph()
        p.text = t
        p.font.size = PPt(23)

    out = SUBMISSIONS / "CBRS503_Presentation_Charbel_Lebbous_8324.pptx"
    prs.save(out)
    return out


def main() -> None:
    report = json.loads(REPORT_JSON.read_text(encoding="utf-8"))
    baseline = json.loads(BASELINE_JSON.read_text(encoding="utf-8"))
    comparison_md = COMPARISON_MD.read_text(encoding="utf-8")
    prompt_text = PROMPT_LOG.read_text(encoding="utf-8")

    versions = parse_prompt_versions(prompt_text)
    cmp_metrics = parse_comparison_metrics(comparison_md)
    charts = build_charts(versions, cmp_metrics)
    report_out = build_report_docx(report, baseline, comparison_md, versions, cmp_metrics, charts)
    ppt_out = build_presentation(versions, cmp_metrics, charts)

    print(f"DOCX: {report_out}")
    print(f"PPTX: {ppt_out}")


if __name__ == "__main__":
    main()
